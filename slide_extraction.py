from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from deplot import process_graph
from pptx.shapes.placeholder import PlaceholderPicture
from openpyxl import load_workbook
from PIL import Image
import io
import argparse

def visitor(shape, slide_idx, image_idx, save_directory):
    try:
        # Attempt to check the shape type
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if isinstance(shape, PlaceholderPicture):
                image_idx = write_image(shape, slide_idx, image_idx, save_directory)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                image_idx = visitor(s, slide_idx, image_idx, save_directory)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_idx = write_image(shape, slide_idx, image_idx, save_directory)
    except NotImplementedError:
        # Log the issue and skip this shape
        print(f"Skipping a shape on slide {slide_idx} of unrecognized type.")
    except AttributeError:
        # This catches the previously handled AttributeError
        print(f"Skipping a shape on slide {slide_idx} that does not have an accessible image.")
    return image_idx


def write_image(shape, slide_idx, image_idx, save_directory):
    try:
        image = shape.image
        image_bytes = image.blob

        # Attempt to get the image extension if available
        try:
            image_ext = image.ext
        except Exception:
            image_ext = 'unknown'

        # Open the image using Pillow
        if image_ext.lower() == 'wmf':
            # Save the WMF image temporarily to convert it
            temp_path = os.path.join(save_directory, f'temp_slide{slide_idx}_image{image_idx:03d}.wmf')
            with open(temp_path, 'wb') as temp_file:
                temp_file.write(image_bytes)

            # Convert WMF to PNG
            image = Image.open(temp_path)
            image_filename = f'slide{slide_idx}_image{image_idx:03d}.png'
            full_path = os.path.join(save_directory, image_filename)
            image.save(full_path, format='PNG')

            # Remove temporary WMF file
            os.remove(temp_path)
        else:
            # Convert other images to PNG
            image = Image.open(io.BytesIO(image_bytes))
            image_filename = f'slide{slide_idx}_image{image_idx:03d}.png'
            full_path = os.path.join(save_directory, image_filename)
            image.save(full_path, format='PNG')

        print("Saving:", full_path)
        return full_path  # Return the full path
    except (AttributeError, KeyError, IOError) as e:
        # Log the issue and skip this shape for any exception
        print(f"Skipping a shape on slide {slide_idx} due to an error: {e}")
        return None  # Return None if there's an error


def get_shape_type_name(shape_type):
    shape_type_map = {
        MSO_SHAPE_TYPE.AUTO_SHAPE: "AutoShape",
        MSO_SHAPE_TYPE.CALLOUT: "Callout",
        MSO_SHAPE_TYPE.CHART: "Chart",
        MSO_SHAPE_TYPE.COMMENT: "Comment",
        MSO_SHAPE_TYPE.FREEFORM: "Freeform",
        MSO_SHAPE_TYPE.GROUP: "Group",
        MSO_SHAPE_TYPE.LINE: "Line",
        MSO_SHAPE_TYPE.MEDIA: "Media",
        MSO_SHAPE_TYPE.PICTURE: "Picture",
        MSO_SHAPE_TYPE.PLACEHOLDER: "Placeholder",
        MSO_SHAPE_TYPE.TABLE: "Table",
        MSO_SHAPE_TYPE.TEXT_BOX: "TextBox",
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: "Embedded OLE Object",
    }
    return shape_type_map.get(shape_type, "Unknown")


def extract_text_from_paragraph(paragraph):
    bullet_point = '-' if paragraph.level >= 0 else '\n'
    left_indent = '  ' * paragraph.level
    return f"{left_indent}{bullet_point}{paragraph.text.strip()}"


def extract_text_from_group_shape(group_shape, level=0):
    text = ""
    for shape in group_shape.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            text += extract_text_from_group_shape(shape, level + 1)
        elif hasattr(shape, "text") and shape.text.strip() != '':
            text += f"{shape.text.strip()}\n"
    return text.strip()


def table_to_markdown(table, index):
    row_contents = []

    for row in table.rows:
        cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
        row_contents.append(cells)

    markdown_table = []
    for row_index, row_data in enumerate(row_contents):
        if row_index == 0:
            header_row = "| " + " | ".join(row_data) + " |"
            separator_row = " | " + " | ".join(['---'] * len(row_data)) + " |"
            markdown_table.append(header_row)
            markdown_table.append(separator_row)
        else:
            content_row = " | " + " | ".join(row_data) + " |"
            markdown_table.append(content_row)

    return f"\n# Table on Slide {index + 1}: \n\n" + "\n".join(markdown_table) + "\n\n -End of Table- \n"


def array_to_markdown(data_2d_array, index):
    row_contents = data_2d_array

    markdown_table = []
    for row_index, row_data in enumerate(row_contents):
        if row_index == 0:
            header_row = "| " + " | ".join(map(str, row_data)) + " |"
            separator_row = "| " + " | ".join(['---'] * len(row_data)) + " |"
            markdown_table.append(header_row)
            markdown_table.append(separator_row)
        else:
            content_row = "| " + " | ".join(map(str, row_data)) + " |"
            markdown_table.append(content_row)

    return f"\n# Table on Slide {index + 1} : \n\n" + "\n".join(markdown_table) + "\n"



def read_pptx_content(pptx_path) -> dict:
    res = {}
    prs = Presentation(pptx_path)

    for slide_number, slide in enumerate(prs.slides):
        slide_content = {"Information": ""}
        information_text = ""
        image_path = None
        count = 1

        for shape in slide.shapes:
            print(shape)
            try:

                if hasattr(shape, "text") and shape.text_frame and shape.text.strip() != '':
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip() != '':
                            information_text += f"\n{extract_text_from_paragraph(paragraph)}"
                    count += 1

                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_text = extract_text_from_group_shape(shape)
                    if group_text:
                        information_text += f"\n{group_text}\n"

                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    information_text += '\n' + table_to_markdown(shape.table, slide_number) + '\n'

                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    chart = shape.chart
                    xlsx_blob = chart.part.chart_workbook.xlsx_part.blob
                    # Load the Excel data into openpyxl
                    xlsx_io = io.BytesIO(xlsx_blob)
                    workbook = load_workbook(xlsx_io)
                    worksheet = workbook.active
                    # Initialize a 2D array to store the data
                    data_2d_array = []
                    # Populate the 2D array with worksheet data
                    for row in worksheet.iter_rows(values_only=True):
                        data_2d_array.append(list(row))

                    markdown = array_to_markdown(data_2d_array, slide_number)
                    information_text += '\n' + markdown + '\n'

                elif shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
                    print(shape.ole_format)

                    ole_format = shape.ole_format
                    print(ole_format.prog_id)
                    if ole_format.prog_id.startswith('Excel.Sheet'):

                        xlsx_blob = ole_format.blob

                        xlsx_io = io.BytesIO(xlsx_blob)
                        workbook = load_workbook(xlsx_io)
                        worksheet = workbook.active
                        data_2d_array = []

                        for row in worksheet.iter_rows(values_only=True):
                            data_2d_array.append(list(row))

                        information_text += '\n' + array_to_markdown(data_2d_array) + '\n'

            except Exception as e:
                try:
                    shape_type_str = get_shape_type_name(shape.shape_type)
                except NotImplementedError:
                    shape_type_str = "Unrecognized Shape Type"
                print(f"Error processing shape on slide {slide_number + 1}, Shape Type: {shape_type_str}: {e}")

        slide_content['Information'] = information_text.strip()
        res[slide_number + 1] = slide_content

    return res


def dict_to_markdown(resulting_dict, md_path):
    with open(md_path, 'w') as md_file:
        for slide_num, content in resulting_dict.items():
            md_file.write(f"# Slide {slide_num}:\n\n")
            information = content['Information'].replace("\n", "  \n")
            md_file.write(f"\n{information}\n\n")
            md_file.write("---\n\n")

def process_directory(input_directory, output_directory):
    for root, dirs, files in os.walk(input_directory):
        for filename in files:
            if filename.endswith('.pptx'):
                file_path = os.path.join(root, filename)
                md_filename = os.path.splitext(filename)[0] + '.md'
                md_path = os.path.join(output_directory, md_filename)

                # Ensure the output directory exists
                os.makedirs(output_directory, exist_ok=True)

                resulting_dict = read_pptx_content(file_path)
                dict_to_markdown(resulting_dict, md_path)

                print(f"Markdown file '{md_path}' has been created for presentation '{filename}'.")

def main():
    input_directory = r'/Users/andrewchung/Downloads/RTE_data/slide_pptx'
    output_directory = '/Users/andrewchung/Downloads/RTE_data/original_MD'
    process_directory(input_directory, output_directory)

if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument("--input_directory", type=str, default='./slide', help='The slides in the directory should b .pptx')
    parser.add_argument("--output_directory", type=str, default='./md')
    args = parser.parse_args()

    process_directory(args.input_directory, args.output_directory)
